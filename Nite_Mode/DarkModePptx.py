from pptx import Presentation
from pptx.slide import Slides
from pptx.dml.color import RGBColor
import pptx.text.text
import os
from django.core.files import File
from django.http import HttpResponse


def convert_pptx(base_dir,filename,filePathName):

    prs = Presentation(filePathName)


    for slide in prs.slides:
        '''Changing background color'''
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        '''Adding Color to the text'''
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            try:
                i = 0
                while(text_frame.paragraphs[i]):
                    p = text_frame.paragraphs[i]
                    p.font.color.rgb = RGBColor(255,255,255)
                    i = i+1
            except:
                pass


    '''Saving the File'''
    path_to_file = os.path.join(base_dir,'media\Main.pptx')
    prs.save(path_to_file)
    print("Han ho gaya ab")
    f = open(path_to_file, 'rb')
    myfile = File(f)
    response = HttpResponse(myfile, content_type='application/[vnd.openxmlformats-officedocument.presentationml.presentation]')
    response['Content-Disposition'] = 'attachment; filename='  + 'DarkFile.pptx'
    return response